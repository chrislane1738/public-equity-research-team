#!/usr/bin/env python
"""Seed a fake DEMO ticker so the UI can be exercised without LLM cost.

Drops one of every previewable file type (md, json, png, xlsx, docx, pdf,
pptx, jsonl) under <RESEARCH_DIR>/DEMO/, mirroring the layout a real
deep-dive would produce. Also inserts a fake "complete" job row into SQLite
so the Sidebar's recent-tickers list surfaces DEMO immediately.

Usage:
    source backend/venv/bin/activate
    python scripts/seed_demo.py            # seeds DEMO
    python scripts/seed_demo.py FAKE       # seeds FAKE instead

Idempotent — running again overwrites existing files in place.
"""
from __future__ import annotations

import asyncio
import json
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

# Ensure the repo root is on sys.path so `backend.*` imports resolve when this
# script is launched directly (not via -m).
_REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_REPO_ROOT))

from backend.config import get_settings  # noqa: E402
from backend.db.sqlite_client import SqliteClient  # noqa: E402


def _write_md(p: Path, ticker: str, kind: str) -> None:
    p.write_text(
        f"# {kind.title()} — {ticker} (demo)\n\n"
        f"This is seed data so the UI can be verified without spending money.\n\n"
        f"## Highlights\n\n"
        f"- Bullet one\n"
        f"- Bullet two\n"
        f"- Bullet three\n\n"
        f"## Numbers\n\n"
        f"| Metric | Value |\n"
        f"|---|---|\n"
        f"| Revenue | $42B |\n"
        f"| EBITDA margin | 38% |\n"
        f"| FCF yield | 4.1% |\n"
    )


def _write_json(p: Path) -> None:
    p.write_text(
        json.dumps(
            {
                "kpis": {
                    "subscribers_m": 287.5,
                    "arpu_usd": 11.62,
                    "data_center_revenue_b": 26.3,
                },
                "notes": "Bespoke operating metrics — seed data.",
            },
            indent=2,
        )
    )


def _write_png(p: Path) -> None:
    fig, ax = plt.subplots(figsize=(6, 4), facecolor="none")
    labels = ["NVDA", "AMD", "INTC", "AVGO", "ARM"]
    shares = [62, 18, 8, 7, 5]
    ax.barh(labels, shares, color="#5b8def")
    ax.set_xlabel("Market share (%)")
    ax.set_title("Demo peer share chart")
    fig.tight_layout()
    fig.savefig(p, dpi=120, transparent=True)
    plt.close(fig)


def _write_xlsx(p: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    cover = wb.active
    cover.title = "Cover"
    cover["A1"] = "DEMO DCF"
    cover["A2"] = "Seed data — not a real valuation"
    rev = wb.create_sheet("Revenue Build")
    rev.append(["Year", "Revenue ($M)", "Growth %"])
    for i, (y, r, g) in enumerate(
        [(2025, 42_000, 12), (2026, 48_500, 15), (2027, 55_400, 14),
         (2028, 62_300, 12), (2029, 68_500, 10)]
    ):
        rev.append([y, r, g])
    sens = wb.create_sheet("Sensitivities")
    sens.append(["WACC \\ g", "2.0%", "2.5%", "3.0%"])
    sens.append(["8.5%", 152, 168, 188])
    sens.append(["9.0%", 138, 151, 167])
    sens.append(["9.5%", 126, 137, 150])
    wb.save(p)


def _write_docx(p: Path, ticker: str) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(f"{ticker} — Investment Memo (demo)", level=1)
    doc.add_paragraph(
        "Executive Summary. This is seed data used to exercise the in-browser "
        "DOCX previewer. The mammoth library converts the file to HTML on the "
        "client and renders it in the workspace center pane."
    )
    doc.add_heading("Investment Thesis", level=2)
    doc.add_paragraph(
        "Bull case: durable competitive moat, share gains, and operating leverage. "
        "Bear case: cyclical demand, regulatory overhang, and cost inflation. "
        "Base case: revenue compounds in the mid-teens with 100bps of annual "
        "margin expansion."
    )
    doc.add_heading("Valuation", level=2)
    doc.add_paragraph(
        "DCF and comps triangulate to a 12-month price target of $185 (+18% upside). "
        "Rating: Buy."
    )
    doc.save(p)


def _write_pdf(p: Path, ticker: str) -> None:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(p), pagesize=LETTER,
                            leftMargin=0.5 * inch, rightMargin=0.5 * inch,
                            topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    flow = [
        Paragraph(f"{ticker} — One-pager (demo)", styles["Title"]),
        Spacer(1, 0.2 * inch),
        Paragraph("Rating: Buy · PT: $185 · Current: $157 (+18%)", styles["Heading2"]),
        Spacer(1, 0.2 * inch),
        Paragraph(
            "Seed data for verifying the react-pdf in-browser previewer. "
            "A real one-pager would contain a thesis, valuation triangulation "
            "table, and top-3 risks.",
            styles["BodyText"],
        ),
    ]
    doc.build(flow)


def _write_pptx(p: Path, ticker: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(layout)
    slide.shapes.title.text = f"{ticker} — Pitch Deck (demo)"
    slide.placeholders[1].text = "Seed data · Buy · PT $185"

    body_layout = pres.slide_layouts[1]
    for title, body in [
        ("Investment Thesis", "Why we like, why now, top risk."),
        ("Business Snapshot", "Segment mix, revenue trend."),
        ("Valuation Triangulation", "DCF · Comps · 52-wk → weighted PT."),
        ("Recommendation", "Buy. Sized 3% of book. Time horizon: 18 months."),
    ]:
        s = pres.slides.add_slide(body_layout)
        s.shapes.title.text = title
        s.placeholders[1].text = body

    pres.save(p)


def _write_jsonl(p: Path, job_id: str) -> None:
    lines = [
        {"ts": datetime.now(timezone.utc).isoformat(), "type": "stage",
         "job_id": job_id, "stage": "stage_1", "status": "started"},
        {"ts": datetime.now(timezone.utc).isoformat(), "type": "agent_completed",
         "job_id": job_id, "agent": "fundamentals",
         "input_tokens": 4200, "output_tokens": 1800, "cost_usd": 0.12,
         "stop_reason": "end_turn"},
        {"ts": datetime.now(timezone.utc).isoformat(), "type": "agent_completed",
         "job_id": job_id, "agent": "dcf",
         "input_tokens": 5100, "output_tokens": 2400, "cost_usd": 0.18,
         "stop_reason": "end_turn"},
        {"ts": datetime.now(timezone.utc).isoformat(), "type": "job_terminal",
         "job_id": job_id, "status": "complete"},
    ]
    p.write_text("\n".join(json.dumps(line) for line in lines) + "\n")


async def _seed_job_row(ticker: str, job_id: str) -> None:
    settings = get_settings()
    sqlite = SqliteClient(settings.sqlite_path)
    await sqlite.connect()
    await sqlite.init_schema()
    # Idempotent: delete any prior demo job for this ticker first.
    await sqlite.execute(
        "DELETE FROM jobs WHERE ticker = ? AND id LIKE 'demo-%'", (ticker,)
    )
    await sqlite.execute(
        "INSERT INTO jobs (id, ticker, workflow, status, current_stage, "
        "agents_status, created_at, completed_at) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (
            job_id,
            ticker,
            "full-deep-dive",
            "complete",
            None,
            json.dumps({"_rating": "Buy"}),
            datetime.now(timezone.utc).isoformat(),
            datetime.now(timezone.utc).isoformat(),
        ),
    )
    await sqlite.close()


def main() -> None:
    ticker = sys.argv[1].upper() if len(sys.argv) > 1 else "DEMO"
    settings = get_settings()
    base = Path(settings.research_dir) / ticker
    job_id = f"demo-{uuid.uuid4().hex[:8]}"

    layout = {
        "fundamentals": ["section.md", "kpis.json"],
        "industry": ["section.md", "peer-share-chart.png"],
        "dcf": ["section.md", "dcf.xlsx"],
        "comps": ["section.md"],
        "macro": ["section.md"],
        "risk": ["section.md"],
        "technicals": ["section.md"],
        "synthesis": ["_synthesis.md"],
        "reports": ["memo.docx", "onepager.pdf", "pitch.pptx"],
        "_logs": [f"{job_id}.jsonl"],
    }

    print(f"Seeding {base}/")
    for folder, files in layout.items():
        d = base / folder
        d.mkdir(parents=True, exist_ok=True)
        for fname in files:
            p = d / fname
            ext = p.suffix.lower()
            print(f"  -> {folder}/{fname}")
            if ext == ".md":
                _write_md(p, ticker, p.stem)
            elif ext == ".json":
                _write_json(p)
            elif ext == ".png":
                _write_png(p)
            elif ext == ".xlsx":
                _write_xlsx(p)
            elif ext == ".docx":
                _write_docx(p, ticker)
            elif ext == ".pdf":
                _write_pdf(p, ticker)
            elif ext == ".pptx":
                _write_pptx(p, ticker)
            elif ext == ".jsonl":
                _write_jsonl(p, job_id)

    asyncio.run(_seed_job_row(ticker, job_id))
    print(f"\nDone. Reload http://localhost:3000 — {ticker} should appear in the sidebar.")


if __name__ == "__main__":
    main()
