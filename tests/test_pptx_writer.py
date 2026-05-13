from pathlib import Path

import pytest
from pptx import Presentation

from backend.tools.pptx_writer import write_pitch_deck, SLIDE_TITLES


@pytest.fixture
def chart_path(tmp_path):
    """Create a tiny valid PNG to use as a placeholder chart."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    p = tmp_path / "chart.png"
    fig, ax = plt.subplots()
    ax.plot([1, 2, 3], [3, 2, 1])
    fig.savefig(p, transparent=True)
    plt.close(fig)
    return p


def test_write_pitch_deck_creates_14_slides(tmp_path, chart_path):
    out = tmp_path / "pitch.pptx"
    write_pitch_deck(
        path=out, ticker="NVDA", rating="Buy",
        price_target=158.0, current_price=110.0,
        slide_bodies={title: f"Body for {title}" for title in SLIDE_TITLES},
        chart_paths={
            "Business Snapshot": chart_path,
            "Industry & Moat": chart_path,
            "Forecast": chart_path,
            "DCF": chart_path,
            "Comps": chart_path,
            "Catalysts": chart_path,
            "Technical Setup": chart_path,
        },
    )
    assert out.exists()
    pres = Presentation(out)
    assert len(pres.slides) == 14


def test_write_pitch_deck_title_slide_contains_ticker_and_pt(tmp_path):
    out = tmp_path / "pitch.pptx"
    write_pitch_deck(
        path=out, ticker="NVDA", rating="Buy",
        price_target=158.0, current_price=110.0,
        slide_bodies={t: "x" for t in SLIDE_TITLES}, chart_paths={},
    )
    pres = Presentation(out)
    title_slide = pres.slides[0]
    text = "\n".join(s.text for s in title_slide.shapes if s.has_text_frame)
    assert "NVDA" in text
    assert "Buy" in text
    assert "158" in text


def test_write_pitch_deck_orders_slides_per_spec(tmp_path):
    out = tmp_path / "pitch.pptx"
    write_pitch_deck(
        path=out, ticker="NVDA", rating="Buy",
        price_target=158.0, current_price=110.0,
        slide_bodies={t: "x" for t in SLIDE_TITLES}, chart_paths={},
    )
    pres = Presentation(out)
    titles = []
    for slide in list(pres.slides)[1:]:  # slide 0 is title
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                titles.append(shape.text_frame.text.split("\n")[0])
                break
    expected_after_title = SLIDE_TITLES[1:]
    assert titles[:len(expected_after_title)] == expected_after_title
