"""python-pptx wrapper that builds a 14-slide pitch deck.

Slides 2-14 use a shared title-on-top + body-on-left + (optional) chart-on-right
layout. The title slide leads with ticker · rating · PT · current price · upside %.
"""
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt


SLIDE_TITLES = [
    "Title",
    "Investment Thesis",
    "Business Snapshot",
    "Industry & Moat",
    "Bespoke KPIs",
    "Financial Performance",
    "Forecast",
    "DCF",
    "Comps",
    "Valuation Triangulation",
    "Catalysts",
    "Risks / Bear Case",
    "Technical Setup",
    "Recommendation",
]


def _add_title_slide(pres: Any, ticker: str, rating: str,
                     price_target: float, current_price: float) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    upside = (price_target - current_price) / current_price * 100 if current_price else 0
    txt = (
        f"{ticker} · {rating}\n"
        f"PT ${price_target:.0f}  ·  Current ${current_price:.0f}  ·  Upside {upside:+.1f}%"
    )
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = box.text_frame
    tf.text = txt
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(36)
            run.font.bold = True


def _add_body_slide(pres: Any, title: str, body: str, chart_path: Path | None) -> None:
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.text = title
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True

    body_w = Inches(5.4) if chart_path else Inches(9)
    body_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), body_w, Inches(5.5))
    body_box.text_frame.word_wrap = True
    body_box.text_frame.text = body
    for p in body_box.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(14)

    if chart_path:
        try:
            slide.shapes.add_picture(str(chart_path), Inches(6.1), Inches(1.4),
                                     width=Inches(3.6))
        except FileNotFoundError:
            # Chart promised by the caller but missing on disk — log via the
            # body box and continue so the rest of the deck still ships.
            note = body_box.text_frame.add_paragraph()
            note.text = f"[chart unavailable: {Path(chart_path).name}]"
            for run in note.runs:
                run.font.italic = True
                run.font.size = Pt(10)


def write_pitch_deck(
    path: Path,
    ticker: str,
    rating: str,
    price_target: float,
    current_price: float,
    slide_bodies: dict[str, str],
    chart_paths: dict[str, Path],
) -> None:
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    pres = Presentation()
    pres.slide_width = Inches(10)
    pres.slide_height = Inches(7.5)

    _add_title_slide(pres, ticker, rating, price_target, current_price)
    for title in SLIDE_TITLES[1:]:
        body = slide_bodies.get(title, "")
        chart = chart_paths.get(title)
        _add_body_slide(pres, title=title, body=body, chart_path=chart)

    pres.save(path)
